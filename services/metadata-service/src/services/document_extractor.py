"""
Document Metadata Extractor

This module provides functionality for extracting metadata from document files,
including PDFs, Office documents, text files, and other document formats.
"""

import os
import json
import asyncio
from typing import Dict, Any, Optional, List, Union
from pathlib import Path
from datetime import datetime
import structlog

from ..core.exceptions import ExtractionError

logger = structlog.get_logger()


class DocumentExtractor:
    """Extracts metadata from document files using multiple methods"""
    
    def __init__(self):
        self.supported_formats = {
            # PDF formats
            'PDF',
            # Microsoft Office formats
            'DOC', 'DOCX', 'XLS', 'XLSX', 'PPT', 'PPTX',
            # OpenDocument formats
            'ODT', 'ODS', 'ODP', 'ODG', 'ODF',
            # Text formats
            'TXT', 'RTF', 'CSV', 'TSV',
            # Markup formats
            'HTML', 'HTM', 'XML', 'XHTML',
            # Markdown formats
            'MD', 'MARKDOWN',
            # E-book formats
            'EPUB', 'MOBI', 'AZW', 'AZW3',
            # Other document formats
            'PS', 'EPS', 'DVI', 'TEX', 'LATEX'
        }
        
        # Common document metadata field mappings
        self.metadata_mapping = {
            'title': 'title',
            'subject': 'subject',
            'author': 'author',
            'creator': 'creator',
            'producer': 'producer',
            'keywords': 'keywords',
            'description': 'description',
            'comments': 'comments',
            'category': 'category',
            'company': 'company',
            'manager': 'manager',
            'creationdate': 'creation_date',
            'modificationdate': 'modification_date',
            'moddate': 'modification_date',
            'trapped': 'trapped',
            'format': 'format',
            'pages': 'page_count',
            'pagecount': 'page_count',
            'words': 'word_count',
            'wordcount': 'word_count',
            'characters': 'character_count',
            'charactercount': 'character_count',
            'lines': 'line_count',
            'linecount': 'line_count',
            'paragraphs': 'paragraph_count',
            'paragraphcount': 'paragraph_count',
            'language': 'language',
            'version': 'version',
            'application': 'application',
            'appversion': 'application_version',
            'security': 'security',
            'encrypted': 'encrypted',
            'template': 'template',
            'lastmodifiedby': 'last_modified_by',
            'lastsavedby': 'last_saved_by',
            'totaltime': 'total_edit_time',
            'edittime': 'total_edit_time',
            'revision': 'revision',
            'revnumber': 'revision_number'
        }
    
    def is_supported_format(self, file_path: str) -> bool:
        """Check if file format is supported for document extraction"""
        extension = Path(file_path).suffix.upper().lstrip('.')
        return extension in self.supported_formats
    
    async def extract_document_metadata(self, file_path: str) -> Dict[str, Any]:
        """
        Extract document metadata from file using multiple methods
        
        Args:
            file_path: Path to document file
            
        Returns:
            Dictionary containing extracted document metadata
        """
        try:
            if not os.path.exists(file_path):
                raise ExtractionError(f"File not found: {file_path}")
            
            if not self.is_supported_format(file_path):
                extension = Path(file_path).suffix
                raise ExtractionError(f"Unsupported file format: {extension}")
            
            # Run extraction in thread pool to avoid blocking
            metadata = await asyncio.get_event_loop().run_in_executor(
                None, self._extract_metadata_sync, file_path
            )
            
            logger.info(
                "document_extraction_completed",
                file_path=file_path,
                has_metadata=bool(metadata.get('processed_metadata', {})),
                page_count=metadata.get('document_info', {}).get('page_count')
            )
            
            return metadata
            
        except Exception as e:
            logger.error(
                "document_extraction_failed",
                file_path=file_path,
                error=str(e)
            )
            raise ExtractionError(f"Failed to extract document metadata: {str(e)}")
    
    def _extract_metadata_sync(self, file_path: str) -> Dict[str, Any]:
        """Synchronous document metadata extraction using multiple methods"""
        try:
            # Initialize result structure
            result = {
                'file_info': {
                    'file_path': file_path,
                    'file_size': os.path.getsize(file_path),
                    'file_name': Path(file_path).name,
                    'file_extension': Path(file_path).suffix.lower(),
                    'extracted_at': datetime.utcnow().isoformat(),
                    'extraction_method': 'multi_method',
                    'extraction_tool': 'document_extractor'
                },
                'raw_metadata': {},
                'processed_metadata': {},
                'document_info': {},
                'content_info': {},
                'security_info': {},
                'extraction_errors': []
            }
            
            # Determine file type and use appropriate extractors
            file_extension = Path(file_path).suffix.lower()
            
            if file_extension == '.pdf':
                result['raw_metadata']['pdf'] = self._extract_pdf_metadata(file_path)
            elif file_extension in ['.doc', '.docx', '.xls', '.xlsx', '.ppt', '.pptx']:
                result['raw_metadata']['office'] = self._extract_office_metadata(file_path)
            elif file_extension in ['.odt', '.ods', '.odp', '.odg', '.odf']:
                result['raw_metadata']['opendocument'] = self._extract_opendocument_metadata(file_path)
            elif file_extension in ['.txt', '.rtf', '.csv', '.md', '.markdown']:
                result['raw_metadata']['text'] = self._extract_text_metadata(file_path)
            elif file_extension in ['.html', '.htm', '.xml', '.xhtml']:
                result['raw_metadata']['markup'] = self._extract_markup_metadata(file_path)
            elif file_extension in ['.epub', '.mobi', '.azw', '.azw3']:
                result['raw_metadata']['ebook'] = self._extract_ebook_metadata(file_path)
            else:
                result['raw_metadata']['generic'] = self._extract_generic_metadata(file_path)
            
            # Process and consolidate metadata
            result['processed_metadata'] = self._process_metadata(result['raw_metadata'])
            result['document_info'] = self._extract_document_info(result['raw_metadata'])
            result['content_info'] = self._extract_content_info(result['raw_metadata'])
            result['security_info'] = self._extract_security_info(result['raw_metadata'])
            
            return result
            
        except Exception as e:
            logger.error("document_sync_extraction_failed", error=str(e))
            raise ExtractionError(f"Document metadata extraction failed: {str(e)}")
    
    def _extract_pdf_metadata(self, file_path: str) -> Optional[Dict[str, Any]]:
        """Extract metadata from PDF files"""
        try:
            import PyPDF2
            
            with open(file_path, 'rb') as file:
                reader = PyPDF2.PdfReader(file)
                
                result = {
                    'page_count': len(reader.pages),
                    'is_encrypted': reader.is_encrypted,
                    'metadata': {}
                }
                
                # Extract metadata
                if reader.metadata:
                    for key, value in reader.metadata.items():
                        # Remove /D: prefix from keys
                        clean_key = key.replace('/D:', '').replace('/', '').lower()
                        result['metadata'][clean_key] = str(value) if value else None
                
                # Extract text from first page for content analysis
                if reader.pages:
                    try:
                        first_page_text = reader.pages[0].extract_text()
                        result['first_page_text'] = first_page_text[:500] if first_page_text else None
                        result['estimated_word_count'] = len(first_page_text.split()) * len(reader.pages) if first_page_text else 0
                    except Exception:
                        result['first_page_text'] = None
                        result['estimated_word_count'] = 0
                
                return result
                
        except ImportError:
            logger.warning("PyPDF2 library not available")
            return None
        except Exception as e:
            logger.warning(f"PDF extraction failed: {str(e)}")
            return None
    
    def _extract_office_metadata(self, file_path: str) -> Optional[Dict[str, Any]]:
        """Extract metadata from Microsoft Office documents"""
        try:
            from python_docx import Document
            from openpyxl import load_workbook
            from pptx import Presentation
            
            file_extension = Path(file_path).suffix.lower()
            
            if file_extension in ['.docx']:
                return self._extract_docx_metadata(file_path)
            elif file_extension in ['.xlsx']:
                return self._extract_xlsx_metadata(file_path)
            elif file_extension in ['.pptx']:
                return self._extract_pptx_metadata(file_path)
            else:
                # Try using olefile for older Office formats
                return self._extract_ole_metadata(file_path)
                
        except ImportError:
            logger.warning("Office document libraries not available")
            return None
        except Exception as e:
            logger.warning(f"Office document extraction failed: {str(e)}")
            return None
    
    def _extract_docx_metadata(self, file_path: str) -> Optional[Dict[str, Any]]:
        """Extract metadata from DOCX files"""
        try:
            from docx import Document
            
            doc = Document(file_path)
            core_props = doc.core_properties
            
            result = {
                'metadata': {
                    'title': core_props.title,
                    'author': core_props.author,
                    'subject': core_props.subject,
                    'keywords': core_props.keywords,
                    'comments': core_props.comments,
                    'category': core_props.category,
                    'created': core_props.created.isoformat() if core_props.created else None,
                    'modified': core_props.modified.isoformat() if core_props.modified else None,
                    'last_modified_by': core_props.last_modified_by,
                    'revision': core_props.revision,
                    'language': core_props.language,
                    'version': core_props.version
                },
                'document_info': {
                    'paragraph_count': len(doc.paragraphs),
                    'table_count': len(doc.tables),
                    'section_count': len(doc.sections)
                }
            }
            
            # Extract text content for analysis
            text_content = '\n'.join([para.text for para in doc.paragraphs])
            result['content_info'] = {
                'text_length': len(text_content),
                'word_count': len(text_content.split()) if text_content else 0,
                'character_count': len(text_content),
                'first_paragraph': doc.paragraphs[0].text if doc.paragraphs else None
            }
            
            return result
            
        except Exception as e:
            logger.warning(f"DOCX extraction failed: {str(e)}")
            return None
    
    def _extract_xlsx_metadata(self, file_path: str) -> Optional[Dict[str, Any]]:
        """Extract metadata from XLSX files"""
        try:
            from openpyxl import load_workbook
            
            workbook = load_workbook(file_path, read_only=True)
            props = workbook.properties
            
            result = {
                'metadata': {
                    'title': props.title,
                    'creator': props.creator,
                    'subject': props.subject,
                    'description': props.description,
                    'keywords': props.keywords,
                    'category': props.category,
                    'created': props.created.isoformat() if props.created else None,
                    'modified': props.modified.isoformat() if props.modified else None,
                    'last_modified_by': props.lastModifiedBy,
                    'revision': props.revision,
                    'version': props.version
                },
                'document_info': {
                    'sheet_count': len(workbook.sheetnames),
                    'sheet_names': workbook.sheetnames
                }
            }
            
            # Analyze first sheet for content info
            if workbook.worksheets:
                first_sheet = workbook.worksheets[0]
                result['content_info'] = {
                    'max_row': first_sheet.max_row,
                    'max_column': first_sheet.max_column,
                    'estimated_cell_count': first_sheet.max_row * first_sheet.max_column
                }
            
            workbook.close()
            return result
            
        except Exception as e:
            logger.warning(f"XLSX extraction failed: {str(e)}")
            return None
    
    def _extract_pptx_metadata(self, file_path: str) -> Optional[Dict[str, Any]]:
        """Extract metadata from PPTX files"""
        try:
            from pptx import Presentation
            
            prs = Presentation(file_path)
            core_props = prs.core_properties
            
            result = {
                'metadata': {
                    'title': core_props.title,
                    'author': core_props.author,
                    'subject': core_props.subject,
                    'keywords': core_props.keywords,
                    'comments': core_props.comments,
                    'category': core_props.category,
                    'created': core_props.created.isoformat() if core_props.created else None,
                    'modified': core_props.modified.isoformat() if core_props.modified else None,
                    'last_modified_by': core_props.last_modified_by,
                    'revision': core_props.revision,
                    'version': core_props.version
                },
                'document_info': {
                    'slide_count': len(prs.slides),
                    'layout_count': len(prs.slide_layouts),
                    'master_count': len(prs.slide_masters)
                }
            }
            
            # Extract text from slides
            slide_texts = []
            for slide in prs.slides:
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, 'text'):
                        slide_text.append(shape.text)
                slide_texts.append(' '.join(slide_text))
            
            all_text = ' '.join(slide_texts)
            result['content_info'] = {
                'text_length': len(all_text),
                'word_count': len(all_text.split()) if all_text else 0,
                'character_count': len(all_text),
                'first_slide_text': slide_texts[0] if slide_texts else None
            }
            
            return result
            
        except Exception as e:
            logger.warning(f"PPTX extraction failed: {str(e)}")
            return None
    
    def _extract_ole_metadata(self, file_path: str) -> Optional[Dict[str, Any]]:
        """Extract metadata from OLE-based Office documents"""
        try:
            import olefile
            
            if not olefile.isOleFile(file_path):
                return None
            
            with olefile.OleFileIO(file_path) as ole:
                result = {
                    'metadata': {},
                    'document_info': {}
                }
                
                # Extract document summary information
                if ole.exists('\\x05DocumentSummaryInformation'):
                    # This would require more complex parsing
                    pass
                
                # Extract summary information
                if ole.exists('\\x05SummaryInformation'):
                    # This would require more complex parsing
                    pass
                
                return result
                
        except ImportError:
            logger.warning("olefile library not available")
            return None
        except Exception as e:
            logger.warning(f"OLE extraction failed: {str(e)}")
            return None
    
    def _extract_opendocument_metadata(self, file_path: str) -> Optional[Dict[str, Any]]:
        """Extract metadata from OpenDocument files"""
        try:
            import zipfile
            import xml.etree.ElementTree as ET
            
            with zipfile.ZipFile(file_path, 'r') as zip_file:
                # Read meta.xml for metadata
                if 'meta.xml' in zip_file.namelist():
                    meta_xml = zip_file.read('meta.xml')
                    root = ET.fromstring(meta_xml)
                    
                    # Define namespaces
                    namespaces = {
                        'meta': 'urn:oasis:names:tc:opendocument:xmlns:meta:1.0',
                        'dc': 'http://purl.org/dc/elements/1.1/',
                        'office': 'urn:oasis:names:tc:opendocument:xmlns:office:1.0'
                    }
                    
                    result = {
                        'metadata': {},
                        'document_info': {}
                    }
                    
                    # Extract Dublin Core metadata
                    dc_fields = ['title', 'creator', 'subject', 'description', 'language', 'date']
                    for field in dc_fields:
                        element = root.find(f'.//dc:{field}', namespaces)
                        if element is not None:
                            result['metadata'][field] = element.text
                    
                    # Extract OpenDocument-specific metadata
                    meta_fields = ['keyword', 'generator', 'initial-creator', 'creation-date', 'editing-cycles']
                    for field in meta_fields:
                        element = root.find(f'.//meta:{field}', namespaces)
                        if element is not None:
                            result['metadata'][field.replace('-', '_')] = element.text
                    
                    # Extract document statistics
                    stats_element = root.find('.//meta:document-statistic', namespaces)
                    if stats_element is not None:
                        for attr, value in stats_element.attrib.items():
                            if attr.startswith('{' + namespaces['meta'] + '}'):
                                clean_attr = attr.replace('{' + namespaces['meta'] + '}', '')
                                result['document_info'][clean_attr] = value
                    
                    return result
                
        except Exception as e:
            logger.warning(f"OpenDocument extraction failed: {str(e)}")
            return None
    
    def _extract_text_metadata(self, file_path: str) -> Optional[Dict[str, Any]]:
        """Extract metadata from text files"""
        try:
            # Read file content
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as file:
                content = file.read()
            
            # Basic text analysis
            lines = content.split('\n')
            words = content.split()
            paragraphs = [p for p in content.split('\n\n') if p.strip()]
            
            result = {
                'content_info': {
                    'line_count': len(lines),
                    'word_count': len(words),
                    'character_count': len(content),
                    'paragraph_count': len(paragraphs),
                    'first_line': lines[0] if lines else None,
                    'last_line': lines[-1] if lines else None
                },
                'document_info': {
                    'encoding': 'utf-8',
                    'file_type': 'text'
                }
            }
            
            # Detect if it's a specific text format
            file_extension = Path(file_path).suffix.lower()
            if file_extension == '.md' or file_extension == '.markdown':
                result['document_info']['file_type'] = 'markdown'
                # Could add markdown-specific analysis here
            elif file_extension == '.csv':
                result['document_info']['file_type'] = 'csv'
                # Basic CSV analysis
                result['content_info']['estimated_rows'] = len(lines)
                if lines:
                    result['content_info']['estimated_columns'] = len(lines[0].split(','))
            
            return result
            
        except Exception as e:
            logger.warning(f"Text extraction failed: {str(e)}")
            return None
    
    def _extract_markup_metadata(self, file_path: str) -> Optional[Dict[str, Any]]:
        """Extract metadata from HTML/XML files"""
        try:
            from bs4 import BeautifulSoup
            
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as file:
                content = file.read()
            
            soup = BeautifulSoup(content, 'html.parser')
            
            result = {
                'metadata': {},
                'document_info': {},
                'content_info': {}
            }
            
            # Extract HTML metadata
            if soup.title:
                result['metadata']['title'] = soup.title.string
            
            # Extract meta tags
            for meta in soup.find_all('meta'):
                name = meta.get('name') or meta.get('property')
                content = meta.get('content')
                if name and content:
                    result['metadata'][name.lower()] = content
            
            # Basic content analysis
            text = soup.get_text()
            result['content_info'] = {
                'text_length': len(text),
                'word_count': len(text.split()) if text else 0,
                'character_count': len(text),
                'tag_count': len(soup.find_all())
            }
            
            return result
            
        except ImportError:
            logger.warning("BeautifulSoup library not available")
            return None
        except Exception as e:
            logger.warning(f"Markup extraction failed: {str(e)}")
            return None
    
    def _extract_ebook_metadata(self, file_path: str) -> Optional[Dict[str, Any]]:
        """Extract metadata from e-book files"""
        try:
            import ebooklib
            from ebooklib import epub
            
            file_extension = Path(file_path).suffix.lower()
            
            if file_extension == '.epub':
                book = epub.read_epub(file_path)
                
                result = {
                    'metadata': {
                        'title': book.get_metadata('DC', 'title')[0][0] if book.get_metadata('DC', 'title') else None,
                        'creator': book.get_metadata('DC', 'creator')[0][0] if book.get_metadata('DC', 'creator') else None,
                        'subject': book.get_metadata('DC', 'subject')[0][0] if book.get_metadata('DC', 'subject') else None,
                        'description': book.get_metadata('DC', 'description')[0][0] if book.get_metadata('DC', 'description') else None,
                        'publisher': book.get_metadata('DC', 'publisher')[0][0] if book.get_metadata('DC', 'publisher') else None,
                        'date': book.get_metadata('DC', 'date')[0][0] if book.get_metadata('DC', 'date') else None,
                        'language': book.get_metadata('DC', 'language')[0][0] if book.get_metadata('DC', 'language') else None,
                        'identifier': book.get_metadata('DC', 'identifier')[0][0] if book.get_metadata('DC', 'identifier') else None,
                        'rights': book.get_metadata('DC', 'rights')[0][0] if book.get_metadata('DC', 'rights') else None
                    },
                    'document_info': {
                        'chapter_count': len(book.get_items_of_type(ebooklib.ITEM_DOCUMENT)),
                        'image_count': len(book.get_items_of_type(ebooklib.ITEM_IMAGE)),
                        'spine_count': len(book.spine)
                    }
                }
                
                return result
            
        except ImportError:
            logger.warning("ebooklib library not available")
            return None
        except Exception as e:
            logger.warning(f"E-book extraction failed: {str(e)}")
            return None
    
    def _extract_generic_metadata(self, file_path: str) -> Optional[Dict[str, Any]]:
        """Extract basic metadata from any file"""
        try:
            stat = os.stat(file_path)
            
            result = {
                'file_info': {
                    'size': stat.st_size,
                    'created': datetime.fromtimestamp(stat.st_ctime).isoformat(),
                    'modified': datetime.fromtimestamp(stat.st_mtime).isoformat(),
                    'accessed': datetime.fromtimestamp(stat.st_atime).isoformat(),
                    'mode': oct(stat.st_mode),
                    'uid': stat.st_uid,
                    'gid': stat.st_gid
                }
            }
            
            return result
            
        except Exception as e:
            logger.warning(f"Generic extraction failed: {str(e)}")
            return None
    
    def _process_metadata(self, raw_metadata: Dict[str, Any]) -> Dict[str, Any]:
        """Process and consolidate metadata from different sources"""
        processed = {}
        
        # Combine metadata from all sources
        all_metadata = {}
        for method_name, method_data in raw_metadata.items():
            if isinstance(method_data, dict):
                if 'metadata' in method_data:
                    all_metadata.update(method_data['metadata'])
                else:
                    all_metadata.update(method_data)
        
        # Map to standardized field names
        for raw_key, raw_value in all_metadata.items():
            if raw_value is None:
                continue
                
            # Normalize key
            normalized_key = str(raw_key).lower().replace(' ', '_').replace('-', '_')
            
            # Use mapping if available
            mapped_key = self.metadata_mapping.get(normalized_key, normalized_key)
            
            # Process value
            if isinstance(raw_value, (list, tuple)):
                processed[mapped_key] = [str(v) for v in raw_value if v is not None]
            else:
                processed[mapped_key] = str(raw_value) if raw_value else None
        
        return processed
    
    def _extract_document_info(self, raw_metadata: Dict[str, Any]) -> Dict[str, Any]:
        """Extract document-specific information"""
        document_info = {}
        
        for method_name, method_data in raw_metadata.items():
            if isinstance(method_data, dict) and 'document_info' in method_data:
                document_info.update(method_data['document_info'])
        
        return document_info
    
    def _extract_content_info(self, raw_metadata: Dict[str, Any]) -> Dict[str, Any]:
        """Extract content-related information"""
        content_info = {}
        
        for method_name, method_data in raw_metadata.items():
            if isinstance(method_data, dict) and 'content_info' in method_data:
                content_info.update(method_data['content_info'])
        
        return content_info
    
    def _extract_security_info(self, raw_metadata: Dict[str, Any]) -> Dict[str, Any]:
        """Extract security-related information"""
        security_info = {}
        
        for method_name, method_data in raw_metadata.items():
            if isinstance(method_data, dict):
                if 'is_encrypted' in method_data:
                    security_info['encrypted'] = method_data['is_encrypted']
                if 'security' in method_data:
                    security_info.update(method_data['security'])
        
        return security_info
    
    async def extract_batch(self, file_paths: List[str]) -> Dict[str, Dict[str, Any]]:
        """Extract metadata from multiple document files"""
        results = {}
        
        for file_path in file_paths:
            try:
                results[file_path] = await self.extract_document_metadata(file_path)
            except Exception as e:
                logger.error(
                    "batch_document_extraction_failed",
                    file_path=file_path,
                    error=str(e)
                )
                results[file_path] = {
                    'error': str(e),
                    'file_info': {
                        'extracted_at': datetime.utcnow().isoformat(),
                        'extraction_tool': 'document_extractor',
                        'file_path': file_path,
                        'success': False
                    }
                }
        
        return results
    
    async def get_document_summary(self, file_path: str) -> Dict[str, Any]:
        """Get a summary of document file information"""
        try:
            metadata = await self.extract_document_metadata(file_path)
            
            # Create summary
            summary = {
                'file_path': file_path,
                'file_name': metadata['file_info']['file_name'],
                'file_size': metadata['file_info']['file_size'],
                'file_extension': metadata['file_info']['file_extension'],
                'title': metadata['processed_metadata'].get('title', 'Unknown'),
                'author': metadata['processed_metadata'].get('author', 'Unknown'),
                'subject': metadata['processed_metadata'].get('subject', 'Unknown'),
                'creation_date': metadata['processed_metadata'].get('creation_date', 'Unknown'),
                'modification_date': metadata['processed_metadata'].get('modification_date', 'Unknown'),
                'page_count': metadata['document_info'].get('page_count', 'Unknown'),
                'word_count': metadata['content_info'].get('word_count', 'Unknown'),
                'character_count': metadata['content_info'].get('character_count', 'Unknown'),
                'has_metadata': bool(metadata['processed_metadata']),
                'is_encrypted': metadata['security_info'].get('encrypted', False),
                'extraction_success': True
            }
            
            return summary
            
        except Exception as e:
            logger.error("document_summary_failed", file_path=file_path, error=str(e))
            return {
                'file_path': file_path,
                'error': str(e),
                'extraction_success': False
            }